"""
PHẦN 2: Slides 8-13 — Deep Dive Features
"""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from slide_helpers import *


def build_slide_08_poi_filter(prs):
    """Slide 8 — Smart POI Filtering"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Feature 2: Spatial + Semantic POI Filtering",
                      "Hybrid search kết hợp PostGIS không gian + pgvector ngữ nghĩa", 8)

    # LEFT — Flow
    flow = [
        ("PostgreSQL + PostGIS", "Lọc theo bán kính từ khách sạn", ACCENT_BLUE),
        ("pgvector Embedding", "Semantic similarity matching", ACCENT_PURPLE),
        ("Hybrid Merge", "Kết hợp spatial + semantic score", ACCENT_GREEN),
        ("Priority Filter", "Locked POI → force include", ACCENT_ORANGE),
    ]
    y = Inches(2.3)
    for title, desc, color in flow:
        add_rounded_rect(slide, Inches(0.8), y, Inches(5.5), Inches(0.95), BG_CARD)
        add_accent_line(slide, Inches(0.8), y, Inches(0.15), color)
        add_text_box(slide, Inches(1.2), y + Inches(0.1), Inches(4.8), Inches(0.35),
                     title, font_size=15, bold=True, color=color)
        add_text_box(slide, Inches(1.2), y + Inches(0.5), Inches(4.8), Inches(0.35),
                     desc, font_size=12, color=GRAY_LIGHT)
        y += Inches(1.1)

    # RIGHT — Filter criteria
    add_rounded_rect(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(4.5), BG_CARD)
    add_text_box(slide, Inches(7.3), Inches(2.5), Inches(5), Inches(0.4),
                 "🎯 Filter Criteria", font_size=18, bold=True, color=ACCENT_BLUE)
    criteria = [
        "📍 Distance — bán kính từ hotel (km)",
        "⭐ Rating — đánh giá từ 3.5+",
        "🏷️ Category — culture, food, nature...",
        "🕐 Opening Hours — lọc theo giờ mở cửa",
        "🔒 Locked POI — user chỉ định → ưu tiên tuyệt đối",
        "🧠 Semantic — vector similarity với sở thích",
    ]
    cy = Inches(3.1)
    for c in criteria:
        add_text_box(slide, Inches(7.3), cy, Inches(5), Inches(0.35),
                     c, font_size=13, color=GRAY_LIGHT)
        cy += Inches(0.42)


def build_slide_09_optimization(prs):
    """Slide 9 — Multi-day Optimization Engine"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Feature 3: AI Route Optimization Engine",
                      "2-Stage pipeline: POI Allocation → CVRPTW Solving bằng OR-Tools", 9)

    # Stage 1
    add_rounded_rect(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(2.0),
                     RGBColor(0x15, 0x1F, 0x34))
    add_text_box(slide, Inches(1.1), Inches(2.4), Inches(5), Inches(0.4),
                 "Stage 1 — POI Allocation", font_size=18, bold=True, color=ACCENT_BLUE)
    s1_items = [
        "Cluster POIs theo vùng địa lý",
        "Phân bổ đều vào N ngày",
        "Đảm bảo locked POI được giữ",
        "Cân bằng workload mỗi ngày",
    ]
    sy = Inches(2.9)
    for item in s1_items:
        add_text_box(slide, Inches(1.3), sy, Inches(4.8), Inches(0.3),
                     f"▸ {item}", font_size=13, color=GRAY_LIGHT)
        sy += Inches(0.32)

    # Arrow
    add_text_box(slide, Inches(3.0), Inches(4.4), Inches(1), Inches(0.5),
                 "⬇", font_size=28, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

    # Stage 2
    add_rounded_rect(slide, Inches(0.8), Inches(4.9), Inches(5.5), Inches(2.0),
                     RGBColor(0x0F, 0x1F, 0x15))
    add_text_box(slide, Inches(1.1), Inches(5.0), Inches(5), Inches(0.4),
                 "Stage 2 — CVRPTW Solver", font_size=18, bold=True, color=ACCENT_GREEN)
    s2_items = [
        "OR-Tools CP-SAT / Routing solver",
        "Minimize: travel time + penalties",
        "Ràng buộc: time window, budget, capacity",
        "Output: ordered route tối ưu cho mỗi ngày",
    ]
    sy = Inches(5.5)
    for item in s2_items:
        add_text_box(slide, Inches(1.3), sy, Inches(4.8), Inches(0.3),
                     f"▸ {item}", font_size=13, color=GRAY_LIGHT)
        sy += Inches(0.32)

    # RIGHT — Constraints summary
    add_rounded_rect(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(4.6), BG_CARD)
    add_text_box(slide, Inches(7.3), Inches(2.5), Inches(5), Inches(0.4),
                 "⚙️ Optimization Constraints", font_size=18, bold=True, color=ACCENT_PURPLE)
    constraints = [
        ("✅ Time Window", "Mỗi POI có giờ mở/đóng cửa"),
        ("✅ Travel Time", "OSRM matrix → thời gian di chuyển thực"),
        ("✅ Visit Duration", "Thời gian tham quan ước tính"),
        ("✅ Budget", "Tổng entrance fee ≤ budget_max"),
        ("✅ Multi-day", "Chia POI ra N ngày, mỗi ngày một route"),
        ("✅ Depot", "Xuất phát & kết thúc tại hotel"),
    ]
    cy = Inches(3.1)
    for title, desc in constraints:
        add_text_box(slide, Inches(7.3), cy, Inches(5), Inches(0.3),
                     title, font_size=14, bold=True, color=ACCENT_GREEN)
        add_text_box(slide, Inches(7.3), cy + Inches(0.28), Inches(5), Inches(0.3),
                     desc, font_size=11, color=GRAY_LIGHT)
        cy += Inches(0.6)


def build_slide_10_solving(prs):
    """Slide 10 — Route Solving Logic"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Mathematical Optimization Flow",
                      "OSRM Matrix → Cost Model → OR-Tools Solver → Optimal Route", 10)

    # Pipeline horizontal
    pipeline = [
        ("OSRM\nMatrix", "Duration +\nDistance", ACCENT_BLUE),
        ("Cost\nMatrix", "Travel cost\n+ Penalties", ACCENT_ORANGE),
        ("Constraint\nModel", "Time windows\nBudget caps", ACCENT_PURPLE),
        ("OR-Tools\nSolver", "CVRPTW\noptimize", ACCENT_GREEN),
        ("Optimal\nRoute", "Ordered POIs\nper day", ACCENT_BLUE),
    ]
    x = Inches(0.5)
    for title, desc, color in pipeline:
        add_rounded_rect(slide, x, Inches(2.5), Inches(2.2), Inches(1.8), BG_CARD)
        add_accent_line(slide, x, Inches(2.5), Inches(2.2), color)
        add_text_box(slide, x + Inches(0.15), Inches(2.7), Inches(1.9), Inches(0.7),
                     title, font_size=15, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), Inches(3.4), Inches(1.9), Inches(0.7),
                     desc, font_size=11, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)
        x += Inches(2.45)

    # Arrows
    for i in range(4):
        ax = Inches(2.6) + Inches(2.45) * i
        add_text_box(slide, ax, Inches(3.0), Inches(0.5), Inches(0.5),
                     "→", font_size=26, color=GRAY_MID, alignment=PP_ALIGN.CENTER)

    # Formula
    add_rounded_rect(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(1.8),
                     RGBColor(0x1A, 0x1A, 0x2E))
    add_text_box(slide, Inches(1.1), Inches(4.9), Inches(5), Inches(0.4),
                 "Objective Function", font_size=16, bold=True, color=ACCENT_BLUE)
    add_text_box(slide, Inches(1.1), Inches(5.4), Inches(10), Inches(0.4),
                 "Minimize:  Σ (travel_time[i][j] × x[i][j])  +  Σ (penalty × dropped[k])",
                 font_size=15, color=WHITE, font_name="Consolas")
    add_text_box(slide, Inches(1.1), Inches(5.9), Inches(10), Inches(0.4),
                 "Subject to:  time_window, capacity, budget, visit_duration, depot_return",
                 font_size=13, color=GRAY_LIGHT, font_name="Consolas")


def build_slide_11_reroute(prs):
    """Slide 11 — Dynamic Re-routing (WOW SLIDE)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Feature 4: Just-in-Time Re-routing",
                      "Tái tối ưu lịch trình ngay khi user lệch kế hoạch", 11)

    # Scenario flow — vertical
    scenario = [
        ("🕐", "User arrives late", "Trễ 45 phút so với kế hoạch", ACCENT_ORANGE),
        ("❌", "Skip attraction", "Bỏ qua Lăng Khải Định (hết giờ mở cửa)", ACCENT_RED),
        ("📍", "GPS → Virtual Depot", "Vị trí hiện tại trở thành điểm xuất phát mới", ACCENT_BLUE),
        ("⚡", "Re-optimize", "Layer 4 solve lại phần còn lại (time_limit=15s)", ACCENT_PURPLE),
        ("✅", "Updated itinerary", "Route mới tối ưu cho thời gian còn lại", ACCENT_GREEN),
    ]

    y = Inches(2.2)
    for emoji, title, desc, color in scenario:
        add_rounded_rect(slide, Inches(0.8), y, Inches(5.5), Inches(0.85), BG_CARD)
        add_text_box(slide, Inches(1.1), y + Inches(0.05), Inches(5), Inches(0.35),
                     f"{emoji}  {title}", font_size=15, bold=True, color=color)
        add_text_box(slide, Inches(1.1), y + Inches(0.42), Inches(5), Inches(0.35),
                     desc, font_size=12, color=GRAY_LIGHT)
        y += Inches(0.95)

    # RIGHT — API details
    add_rounded_rect(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(4.8), BG_CARD)
    add_text_box(slide, Inches(7.3), Inches(2.4), Inches(5), Inches(0.4),
                 "🔌 API Flow", font_size=16, bold=True, color=ACCENT_BLUE)

    api_flow = [
        "Mobile → GPS + remaining_poi_ids",
        "POST /v1/trip/re_route (Gateway)",
        "Gateway → build ReRouteRequest",
        "POST /re-route (Layer 4)",
        "Layer 4 → Virtual Depot + solve_day()",
        "Response → TravelItineraryDay",
        "Mobile → merge day + update UI",
    ]
    ay = Inches(3.0)
    for i, step in enumerate(api_flow):
        color = ACCENT_GREEN if i == len(api_flow) - 1 else GRAY_LIGHT
        add_text_box(slide, Inches(7.3), ay, Inches(5), Inches(0.3),
                     f"{i+1}. {step}", font_size=12, color=color)
        ay += Inches(0.38)


def build_slide_12_realtime(prs):
    """Slide 12 — Real-time UX"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Feature 5: Real-time User Experience",
                      "SSE streaming + Push notification + Background tracking", 12)

    cards = [
        ("📡 SSE Streaming", [
            "Server-Sent Events từ Gateway",
            "Event: l2_done → l3_done → l4_result",
            "Hook: useTripPipeline()",
            "Cập nhật UI realtime trên Loading screen",
        ], ACCENT_BLUE),
        ("🔔 Push Notification", [
            "expo-notifications local trigger",
            '"Hành trình đã sẵn sàng!"',
            "Trigger khi pipeline L4 hoàn thành",
            "Auto-navigate sang MapTimeline",
        ], ACCENT_GREEN),
        ("📍 Background Location", [
            "expo-location background task",
            "Track mỗi 100m di chuyển",
            "Nền tảng cho auto-trigger reroute",
            "Geofence detection (future)",
        ], ACCENT_PURPLE),
    ]

    x = Inches(0.5)
    for title, items, color in cards:
        add_rounded_rect(slide, x, Inches(2.3), Inches(3.9), Inches(4.2), BG_CARD)
        add_accent_line(slide, x, Inches(2.3), Inches(3.9), color)
        add_text_box(slide, x + Inches(0.3), Inches(2.5), Inches(3.3), Inches(0.5),
                     title, font_size=17, bold=True, color=color)
        iy = Inches(3.2)
        for item in items:
            add_text_box(slide, x + Inches(0.3), iy, Inches(3.3), Inches(0.35),
                         f"• {item}", font_size=12, color=GRAY_LIGHT)
            iy += Inches(0.38)
        x += Inches(4.15)


def build_slide_13_ui(prs):
    """Slide 13 — Visual Itinerary UI"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Smart MapTimeline Interface",
                      "Mapbox route + Day tabs + Summary bar + Cost estimation", 13)

    # UI Feature cards
    ui_features = [
        ("🗺️ Mapbox Map", "Route GeoJSON\nPOI markers\nUser location", ACCENT_BLUE),
        ("📅 Day Tabs", "Ngày 1 / Ngày 2 / Ngày 3\nSố điểm + chi phí/ngày\nFilter timeline theo ngày", ACCENT_GREEN),
        ("📊 Summary Bar", "Tổng stops · km · thời gian\nEst. Cost tổng\nFlyTo camera animation", ACCENT_PURPLE),
        ("🔄 Re-route Sheet", "Chọn POI muốn bỏ\nXác nhận re-route\nBottomSheet gesture", ACCENT_ORANGE),
    ]

    x = Inches(0.5)
    for title, desc, color in ui_features:
        add_rounded_rect(slide, x, Inches(2.3), Inches(2.9), Inches(3.0), BG_CARD)
        add_accent_line(slide, x, Inches(2.3), Inches(2.9), color)
        add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(2.5), Inches(0.5),
                     title, font_size=16, bold=True, color=color)
        ty = Inches(3.1)
        for line in desc.split("\n"):
            add_text_box(slide, x + Inches(0.2), ty, Inches(2.5), Inches(0.3),
                         f"• {line}", font_size=12, color=GRAY_LIGHT)
            ty += Inches(0.35)
        x += Inches(3.15)

    # Tech note
    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8),
                 "Tech: Expo React Native + TypeScript  •  @rnmapbox/maps  •  @gorhom/bottom-sheet  •  react-native-reanimated",
                 font_size=13, color=GRAY_MID, alignment=PP_ALIGN.CENTER)
