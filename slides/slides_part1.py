"""
PHẦN 1: Slides 1-7 — Mở đầu + Tổng quan giải pháp
"""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from slide_helpers import *


def build_slide_01_cover(prs):
    """Slide 1 — Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Accent lines decorative
    add_accent_line(slide, Inches(0.8), Inches(1.8), Inches(3), ACCENT_BLUE)
    add_accent_line(slide, Inches(0.8), Inches(1.9), Inches(1.5), ACCENT_GREEN)

    # Main title
    add_text_box(slide, Inches(0.8), Inches(2.1), Inches(11), Inches(1.2),
                 "AI-Driven Dynamic\nItinerary Optimizer",
                 font_size=44, bold=True, color=WHITE)

    # Vietnamese subtitle
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(11), Inches(0.6),
                 "Hệ thống tối ưu lịch trình du lịch động bằng AI",
                 font_size=20, color=GRAY_LIGHT)

    # Tech tags
    add_text_box(slide, Inches(0.8), Inches(4.3), Inches(11), Inches(0.5),
                 "NLP  •  Semantic Search  •  Multi-day Route Optimization  •  Real-time Re-routing",
                 font_size=14, color=ACCENT_BLUE)

    # Bottom info
    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(5), Inches(0.4),
                 "GVHD: [Tên GVHD]", font_size=13, color=GRAY_MID)
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(5), Inches(0.4),
                 "Team: [Tên nhóm]  •  Trường: [Tên trường]",
                 font_size=13, color=GRAY_MID)

    # Right side decorative card
    add_rounded_rect(slide, Inches(8.5), Inches(2.5), Inches(4), Inches(3.5),
                     RGBColor(0x15, 0x1F, 0x34))
    add_text_box(slide, Inches(8.8), Inches(2.8), Inches(3.5), Inches(0.4),
                 "📍 Huế, Vietnam", font_size=14, color=ACCENT_BLUE, bold=True)
    items = ["Đại Nội Huế  ·  08:00–10:00",
             "Chùa Thiên Mụ  ·  10:30–11:30",
             "Sông Hương  ·  14:00–16:00",
             "Chợ Đông Ba  ·  16:30–17:30"]
    y = Inches(3.3)
    for item in items:
        add_text_box(slide, Inches(8.8), y, Inches(3.5), Inches(0.3),
                     f"▸ {item}", font_size=11, color=GRAY_LIGHT)
        y += Inches(0.35)
    add_text_box(slide, Inches(8.8), y + Inches(0.1), Inches(3.5), Inches(0.3),
                 "Total: 24.5 km  ·  350k₫  ·  3 days",
                 font_size=12, color=ACCENT_GREEN, bold=True)


def build_slide_02_problem(prs):
    """Slide 2 — Bối cảnh & Vấn đề"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Why Travel Planning is Hard?",
                      "Người dùng đang gặp quá nhiều rào cản khi tự lên lịch trình", 2)

    # LEFT column — Pain points
    pain_points = [
        ("❌  Quá nhiều app", "Google Maps + Blogs + Booking + Weather + TikTok"),
        ("❌  Lịch trình thủ công", "Di chuyển vòng vèo, trùng khu vực, không cân đối thời gian"),
        ("❌  Không thích ứng realtime", "Trễ lịch, huỷ điểm đến, thời tiết xấu → phải làm lại từ đầu"),
    ]
    y = Inches(2.2)
    for title, desc in pain_points:
        add_rounded_rect(slide, Inches(0.8), y, Inches(5.5), Inches(1.2))
        add_text_box(slide, Inches(1.1), y + Inches(0.15), Inches(5), Inches(0.4),
                     title, font_size=18, bold=True, color=ACCENT_RED)
        add_text_box(slide, Inches(1.1), y + Inches(0.6), Inches(5), Inches(0.5),
                     desc, font_size=13, color=GRAY_LIGHT)
        y += Inches(1.4)

    # RIGHT column — Summary
    add_rounded_rect(slide, Inches(7), y_r := Inches(2.2), Inches(5.5), Inches(4.0),
                     RGBColor(0x1E, 0x0F, 0x0F))
    add_text_box(slide, Inches(7.3), y_r + Inches(0.2), Inches(5), Inches(0.5),
                 "Hệ quả", font_size=20, bold=True, color=ACCENT_ORANGE)
    consequences = [
        "⏱️  Mất hàng giờ điều phối thủ công",
        "😤  Trải nghiệm kém liền mạch",
        "💸  Bỏ lỡ POI quan trọng hoặc vượt ngân sách",
        "🔄  Không có cơ chế cập nhật route nhanh",
    ]
    cy = y_r + Inches(0.8)
    for c in consequences:
        add_text_box(slide, Inches(7.3), cy, Inches(5), Inches(0.35),
                     c, font_size=14, color=GRAY_LIGHT)
        cy += Inches(0.45)

    # Punchline
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(12), Inches(0.5),
                 "→ Existing itinerary tools are static, not adaptive.",
                 font_size=16, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)


def build_slide_03_gap(prs):
    """Slide 3 — Khoảng trống thị trường"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Existing Solutions — Limitation",
                      "So sánh các giải pháp hiện có với hệ thống đề xuất", 3)

    # Table
    rows, cols = 7, 6
    tbl_shape = slide.shapes.add_table(rows, cols,
                                        Inches(0.8), Inches(2.2),
                                        Inches(11.5), Inches(3.8))
    tbl = tbl_shape.table

    headers = ["Feature", "Google Maps", "TripAdvisor", "Layla AI", "Traveloka", "Our System"]
    data = [
        ["NL Planning",      "❌", "❌", "✅", "❌", "✅"],
        ["Route Optimization","⚠️", "❌", "⚠️", "❌", "✅"],
        ["Multi-day Plan",   "❌", "❌", "⚠️", "❌", "✅"],
        ["Dynamic Re-route", "❌", "❌", "❌", "❌", "✅"],
        ["Semantic Search",  "❌", "❌", "⚠️", "❌", "✅"],
        ["Budget Constraint","❌", "❌", "❌", "⚠️", "✅"],
    ]

    # Style header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)

    # Data rows
    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.alignment = PP_ALIGN.CENTER
                if j == 5 and val == "✅":
                    p.font.color.rgb = ACCENT_GREEN
                    p.font.bold = True
                elif val == "❌":
                    p.font.color.rgb = ACCENT_RED
                elif val == "⚠️":
                    p.font.color.rgb = ACCENT_ORANGE
                else:
                    p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD if i % 2 == 0 else BG_DARK

    # Punchline
    tf = add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.7),
                      "Existing systems generate suggestions.",
                      font_size=18, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, "We optimize executable itineraries.",
                  font_size=22, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)


def build_slide_04_objectives(prs):
    """Slide 4 — Mục tiêu đề tài"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Research Objectives",
                      "Từ yêu cầu mơ hồ đến lịch trình tối ưu, thực thi được", 4)

    # Funnel flow
    funnel_steps = [
        ("📝", "Input", "Messy travel request", ACCENT_ORANGE),
        ("🤖", "AI Extraction", "NLP intent parsing", ACCENT_BLUE),
        ("⚡", "Optimization", "CVRPTW multi-day solving", ACCENT_PURPLE),
        ("🗺️", "Output", "Smart executable itinerary", ACCENT_GREEN),
    ]
    x = Inches(0.8)
    for emoji, title, desc, color in funnel_steps:
        add_rounded_rect(slide, x, Inches(2.3), Inches(2.7), Inches(1.8), BG_CARD)
        add_text_box(slide, x + Inches(0.2), Inches(2.4), Inches(2.3), Inches(0.5),
                     f"{emoji}  {title}", font_size=18, bold=True, color=color)
        add_text_box(slide, x + Inches(0.2), Inches(3.0), Inches(2.3), Inches(0.8),
                     desc, font_size=13, color=GRAY_LIGHT)
        x += Inches(3.0)

    # Arrows between cards
    for i in range(3):
        ax = Inches(3.4) + Inches(3.0) * i
        add_text_box(slide, ax, Inches(2.9), Inches(0.5), Inches(0.5),
                     "→", font_size=28, color=GRAY_MID, alignment=PP_ALIGN.CENTER)

    # Objectives list
    objectives = [
        "1. Tự động tạo lịch trình du lịch từ ngôn ngữ tự nhiên (Tiếng Việt)",
        "2. Tối ưu đa ngày với ràng buộc thực tế (thời gian, ngân sách, khoảng cách)",
        "3. Hỗ trợ tái tối ưu động (JIT Re-routing) khi phát sinh sự cố",
        "4. Cải thiện trải nghiệm du lịch cá nhân hóa qua Semantic Search",
    ]
    y = Inches(4.5)
    for obj in objectives:
        add_text_box(slide, Inches(1.2), y, Inches(10.5), Inches(0.4),
                     obj, font_size=15, color=WHITE)
        y += Inches(0.45)


def build_slide_05_overview(prs):
    """Slide 5 — Tổng quan hệ thống (KEY SLIDE)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Proposed System Overview",
                      "End-to-end pipeline: từ prompt đến lịch trình tối ưu trên bản đồ", 5)

    # Pipeline flow blocks
    steps = [
        ("User Prompt",        "NL input",                ACCENT_ORANGE, Inches(0.5)),
        ("NLP Extraction",     "LLMExtractorService",     ACCENT_BLUE,   Inches(2.2)),
        ("Semantic Search",    "pgvector + embedding",    ACCENT_PURPLE, Inches(3.9)),
        ("Spatial Filter",     "PostGIS radius + rating", ACCENT_BLUE,   Inches(5.6)),
        ("CVRPTW Solver",      "OR-Tools optimization",   ACCENT_GREEN,  Inches(7.3)),
        ("MapTimeline UI",     "Mapbox + Day tabs",       ACCENT_BLUE,   Inches(9.0)),
        ("JIT Re-route",       "GPS → Virtual Depot",     ACCENT_RED,    Inches(10.7)),
    ]

    for title, desc, color, x_pos in steps:
        add_rounded_rect(slide, x_pos, Inches(2.5), Inches(1.6), Inches(2.0), BG_CARD)
        add_accent_line(slide, x_pos, Inches(2.5), Inches(1.6), color)
        add_text_box(slide, x_pos + Inches(0.1), Inches(2.7), Inches(1.4), Inches(0.7),
                     title, font_size=13, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x_pos + Inches(0.1), Inches(3.5), Inches(1.4), Inches(0.8),
                     desc, font_size=10, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

    # Arrows
    for i in range(6):
        ax = Inches(2.05) + Inches(1.7) * i
        add_text_box(slide, ax, Inches(3.1), Inches(0.4), Inches(0.4),
                     "→", font_size=22, color=GRAY_MID, alignment=PP_ALIGN.CENTER)

    # Bottom labels
    labels = [
        ("Layer 5 — Mobile", Inches(0.5), ACCENT_ORANGE),
        ("Layer 2-3 — Gateway API", Inches(2.2), ACCENT_BLUE),
        ("Layer 4 — Optimization Engine", Inches(7.3), ACCENT_GREEN),
        ("Layer 5 — Mobile UI", Inches(9.0), ACCENT_BLUE),
    ]
    for txt, x, color in labels:
        add_text_box(slide, x, Inches(4.8), Inches(3), Inches(0.3),
                     txt, font_size=11, color=color, alignment=PP_ALIGN.CENTER)


def build_slide_06_users(prs):
    """Slide 6 — Target Users"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Target Users", "Ai sẽ sử dụng hệ thống này?", 6)

    personas = [
        ("👨‍👩‍👧", "Casual Travelers", "Không muốn tự lên lịch\nThích có sẵn itinerary tối ưu\nĐi gia đình, nhóm bạn", ACCENT_BLUE),
        ("💸", "Budget Travelers", "Muốn tối ưu chi phí\nSo sánh entrance fee\nKiểm soát ngân sách từng ngày", ACCENT_GREEN),
        ("🧳", "Smart Independents", "Muốn itinerary cá nhân hóa\nCần realtime re-route\nDu lịch tự túc chuyên nghiệp", ACCENT_PURPLE),
        ("🏃", "Time-constrained", "Du lịch ngắn ngày (1-3 ngày)\nCần tối ưu mỗi phút\nKhông muốn lãng phí thời gian", ACCENT_ORANGE),
    ]

    x = Inches(0.5)
    for emoji, name, desc, color in personas:
        add_rounded_rect(slide, x, Inches(2.3), Inches(2.9), Inches(4.0), BG_CARD)
        add_accent_line(slide, x, Inches(2.3), Inches(2.9), color)
        add_text_box(slide, x + Inches(0.3), Inches(2.6), Inches(2.3), Inches(0.5),
                     f"{emoji}  {name}", font_size=17, bold=True, color=color)
        # Description lines
        y = Inches(3.3)
        for line in desc.split("\n"):
            add_text_box(slide, x + Inches(0.3), y, Inches(2.3), Inches(0.35),
                         f"• {line}", font_size=12, color=GRAY_LIGHT)
            y += Inches(0.35)
        x += Inches(3.15)


def build_slide_07_nlp(prs):
    """Slide 7 — NLP Itinerary Generation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Feature 1: Natural Language → Structured Plan",
                      "LLMExtractorService trích xuất intent thành LLMDataContract", 7)

    # LEFT — Prompt example
    add_rounded_rect(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(1.5),
                     RGBColor(0x1A, 0x1A, 0x2E))
    add_text_box(slide, Inches(1.0), Inches(2.3), Inches(5), Inches(0.3),
                 "💬 User Prompt", font_size=14, bold=True, color=ACCENT_ORANGE)
    add_text_box(slide, Inches(1.0), Inches(2.7), Inches(5), Inches(0.8),
                 '"Tôi muốn đi Đại Nội và dạo bờ sông Hương,\n3 ngày, ngân sách 1 triệu"',
                 font_size=15, color=WHITE)

    # Arrow
    add_text_box(slide, Inches(3.0), Inches(3.8), Inches(1), Inches(0.5),
                 "⬇", font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # LEFT — JSON output
    add_rounded_rect(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(2.5),
                     RGBColor(0x0D, 0x1B, 0x0D))
    add_text_box(slide, Inches(1.0), Inches(4.4), Inches(5), Inches(0.3),
                 "📋 LLMDataContract (Output)", font_size=14, bold=True, color=ACCENT_GREEN)
    json_text = ('{\n'
                 '  "num_days": 3,\n'
                 '  "budget_max": 1000000,\n'
                 '  "locked_pois": ["Đại Nội", "Sông Hương"],\n'
                 '  "tags": ["culture", "relax", "scenic"],\n'
                 '  "preferred_transport": "motorbike"\n'
                 '}')
    add_text_box(slide, Inches(1.0), Inches(4.8), Inches(5), Inches(1.8),
                 json_text, font_size=12, color=ACCENT_GREEN, font_name="Consolas")

    # RIGHT — Flow
    flow_steps = [
        ("1. User Prompt", "Ngôn ngữ tự nhiên (Tiếng Việt)", ACCENT_ORANGE),
        ("2. LLM Extractor", "Gemini API → JSON contract", ACCENT_BLUE),
        ("3. Embedding", "Vector hóa sở thích user", ACCENT_PURPLE),
        ("4. Spatial Filter", "PostGIS + pgvector hybrid search", ACCENT_BLUE),
        ("5. Layer 4 Solver", "OR-Tools CVRPTW optimization", ACCENT_GREEN),
    ]
    y = Inches(2.2)
    for title, desc, color in flow_steps:
        add_rounded_rect(slide, Inches(7), y, Inches(5.5), Inches(0.85), BG_CARD)
        add_text_box(slide, Inches(7.2), y + Inches(0.05), Inches(5), Inches(0.35),
                     title, font_size=14, bold=True, color=color)
        add_text_box(slide, Inches(7.2), y + Inches(0.4), Inches(5), Inches(0.35),
                     desc, font_size=11, color=GRAY_LIGHT)
        y += Inches(0.95)
