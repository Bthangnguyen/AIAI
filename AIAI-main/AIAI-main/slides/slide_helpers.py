"""
Shared helpers: colors, fonts, utility functions for slide generation.
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Color Palette (Dark Premium Theme) ───
BG_DARK       = RGBColor(0x0F, 0x17, 0x2A)  # Deep navy
BG_CARD       = RGBColor(0x1A, 0x25, 0x3C)  # Card background
ACCENT_BLUE   = RGBColor(0x38, 0xBD, 0xF8)  # Bright cyan
ACCENT_GREEN  = RGBColor(0x4A, 0xDE, 0x80)  # Green
ACCENT_ORANGE = RGBColor(0xFB, 0x92, 0x3C)  # Orange
ACCENT_RED    = RGBColor(0xF8, 0x71, 0x71)  # Red/pink
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)  # Purple
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT    = RGBColor(0x94, 0xA3, 0xB8)  # Muted text
GRAY_MID      = RGBColor(0x64, 0x74, 0x8B)
TRANSPARENT   = RGBColor(0x0F, 0x17, 0x2A)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    """Add a simple text box and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(text_frame, text, font_size=16, color=WHITE,
                  bold=False, space_before=Pt(4), space_after=Pt(2),
                  alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Append a paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    p.alignment = alignment
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD):
    """Add a rounded rectangle shape as a card background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    """Add a thin horizontal accent line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def slide_title_block(slide, title, subtitle=None, slide_num=None):
    """Add a consistent title + subtitle block at the top of a content slide."""
    set_slide_bg(slide)
    # Accent line
    add_accent_line(slide, Inches(0.8), Inches(0.5), Inches(2.5))
    # Slide number
    if slide_num:
        add_text_box(slide, Inches(11.5), Inches(0.3), Inches(1.5), Inches(0.4),
                     f"{slide_num:02d}", font_size=12, color=GRAY_MID,
                     alignment=PP_ALIGN.RIGHT)
    # Title
    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
                 title, font_size=32, bold=True, color=WHITE)
    # Subtitle
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
                     subtitle, font_size=16, color=GRAY_LIGHT)


def add_flow_arrow(slide, left, top, width=Inches(0.5), color=ACCENT_BLUE):
    """Add a downward arrow shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                   left, top, width, Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card_with_text(slide, left, top, width, height, title, body_lines,
                       accent_color=ACCENT_BLUE, title_size=16, body_size=13):
    """Add a card (rounded rect) with title and bullet lines inside."""
    add_rounded_rect(slide, left, top, width, height)
    # Card title
    tf = add_text_box(slide, left + Inches(0.2), top + Inches(0.15),
                      width - Inches(0.4), Inches(0.4),
                      title, font_size=title_size, bold=True, color=accent_color)
    # Body lines
    y = top + Inches(0.55)
    for line in body_lines:
        add_text_box(slide, left + Inches(0.25), y,
                     width - Inches(0.5), Inches(0.3),
                     line, font_size=body_size, color=GRAY_LIGHT)
        y += Inches(0.28)
