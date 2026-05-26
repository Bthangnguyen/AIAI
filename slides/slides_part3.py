"""
PHẦN 3: Slides 14-22 — Architecture, Tech Stack, Demo, Evaluation, Closing
"""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from slide_helpers import *


def build_slide_14_architecture(prs):
    """Slide 14 — System Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "System Architecture",
                      "Kiến trúc phân lớp 5-Layer với separation of concerns", 14)

    layers = [
        ("Layer 5 — Mobile App", "Expo React Native + TypeScript\nMapbox, BottomSheet, SSE client\nTripService API calls", ACCENT_ORANGE, Inches(2.2)),
        ("Layer 2-3 — Gateway API", "FastAPI (Python)\nLLMExtractor, SpatialFilter, EmbeddingService\nSSE streaming + rate limiting (slowapi)", ACCENT_BLUE, Inches(3.5)),
        ("Layer 4 — Optimization", "FastAPI (Python)\nTravelPlanService: Allocator + Solver\nOR-Tools CVRPTW engine", ACCENT_GREEN, Inches(4.8)),
        ("Infrastructure", "OSRM (routing matrix) + Haversine fallback\nPostgreSQL + PostGIS + pgvector\nSQLite distance cache • Docker", ACCENT_PURPLE, Inches(6.1)),
    ]

    for title, desc, color, y in layers:
        add_rounded_rect(slide, Inches(0.8), y, Inches(11.5), Inches(1.1), BG_CARD)
        add_accent_line(slide, Inches(0.8), y, Inches(0.2), color)
        add_text_box(slide, Inches(1.2), y + Inches(0.1), Inches(3.5), Inches(0.4),
                     title, font_size=16, bold=True, color=color)
        lines = desc.split("\n")
        lx = Inches(5.0)
        for line in lines:
            add_text_box(slide, lx, y + Inches(0.1), Inches(3.5), Inches(0.35),
                         f"• {line}", font_size=11, color=GRAY_LIGHT)
            lx += Inches(0.0)
        # Put desc as single block
        add_text_box(slide, Inches(5.0), y + Inches(0.15), Inches(7), Inches(0.9),
                     desc.replace("\n", "  •  "), font_size=11, color=GRAY_LIGHT)

    # Arrows between layers
    for y_pos in [Inches(3.35), Inches(4.65), Inches(5.95)]:
        add_text_box(slide, Inches(6.0), y_pos, Inches(0.5), Inches(0.3),
                     "⬇", font_size=16, color=GRAY_MID, alignment=PP_ALIGN.CENTER)


def build_slide_15_pipeline(prs):
    """Slide 15 — Data Flow Pipeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "End-to-End Data Pipeline",
                      "Sequence: Request → Extract → Search → Optimize → Render", 15)

    steps = [
        ("1", "User Request", "prompt + hotel + days", ACCENT_ORANGE),
        ("2", "LLM Extraction", "intent → LLMDataContract", ACCENT_BLUE),
        ("3", "Embedding", "vector hóa preferences", ACCENT_PURPLE),
        ("4", "Spatial Query", "PostGIS + pgvector", ACCENT_BLUE),
        ("5", "Build Payload", "Layer 4 format", GRAY_LIGHT),
        ("6", "Optimization", "OR-Tools CVRPTW solve", ACCENT_GREEN),
        ("7", "SSE Stream", "progress events", ACCENT_ORANGE),
        ("8", "UI Rendering", "MapTimeline display", ACCENT_BLUE),
    ]

    x = Inches(0.3)
    for num, title, desc, color in steps:
        w = Inches(1.5)
        add_rounded_rect(slide, x, Inches(2.8), w, Inches(2.5), BG_CARD)
        add_accent_line(slide, x, Inches(2.8), w, color)
        # Number circle
        add_text_box(slide, x + Inches(0.1), Inches(2.95), Inches(0.4), Inches(0.4),
                     num, font_size=18, bold=True, color=color)
        add_text_box(slide, x + Inches(0.1), Inches(3.4), w - Inches(0.2), Inches(0.6),
                     title, font_size=12, bold=True, color=WHITE)
        add_text_box(slide, x + Inches(0.1), Inches(4.0), w - Inches(0.2), Inches(0.8),
                     desc, font_size=10, color=GRAY_LIGHT)
        x += Inches(1.6)

    # API endpoints reference
    add_rounded_rect(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0),
                     RGBColor(0x1A, 0x1A, 0x2E))
    add_text_box(slide, Inches(1.1), Inches(5.9), Inches(11), Inches(0.35),
                 "API Endpoints (verified from codebase):",
                 font_size=13, bold=True, color=ACCENT_BLUE)
    add_text_box(slide, Inches(1.1), Inches(6.3), Inches(11), Inches(0.35),
                 "POST /v1/trip/plan_trip_stream  •  POST /v1/trip/re_route  •  POST /plan  •  POST /re-route  •  GET /health",
                 font_size=12, color=ACCENT_GREEN, font_name="Consolas")


def build_slide_16_techstack(prs):
    """Slide 16 — Tech Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Technology Stack",
                      "Lý do lựa chọn từng công nghệ trong hệ thống", 16)

    categories = [
        ("Backend", [
            ("FastAPI", "API nhanh, async, rõ contract"),
            ("PostgreSQL", "Robust RDBMS cho POI data"),
        ], ACCENT_BLUE),
        ("AI / Search", [
            ("Gemini LLM", "Intent extraction (NLP)"),
            ("pgvector", "Semantic similarity search"),
        ], ACCENT_PURPLE),
        ("Optimization", [
            ("OR-Tools", "CVRPTW solver mạnh mẽ"),
            ("OSRM", "Real-world travel time matrix"),
        ], ACCENT_GREEN),
        ("Frontend", [
            ("React Native", "Cross-platform mobile"),
            ("Mapbox GL", "Interactive map rendering"),
        ], ACCENT_ORANGE),
        ("Infra", [
            ("Docker", "Container deployment"),
            ("SSE + Notification", "Realtime experience"),
        ], GRAY_LIGHT),
    ]

    x = Inches(0.3)
    for cat_name, items, color in categories:
        add_rounded_rect(slide, x, Inches(2.3), Inches(2.4), Inches(4.2), BG_CARD)
        add_accent_line(slide, x, Inches(2.3), Inches(2.4), color)
        add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(2), Inches(0.4),
                     cat_name, font_size=16, bold=True, color=color)
        iy = Inches(3.1)
        for tech, reason in items:
            add_text_box(slide, x + Inches(0.2), iy, Inches(2), Inches(0.35),
                         tech, font_size=14, bold=True, color=WHITE)
            add_text_box(slide, x + Inches(0.2), iy + Inches(0.32), Inches(2), Inches(0.6),
                         reason, font_size=11, color=GRAY_LIGHT)
            iy += Inches(0.9)
        x += Inches(2.55)


def build_slide_17_fallback(prs):
    """Slide 17 — Reliability & Fallback"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Fault Tolerance & Fallback Strategy",
                      "Hệ thống luôn trả kết quả — kể cả khi có lỗi", 17)

    fallbacks = [
        ("🗺️ OSRM Failed", "→ Haversine fallback", "Độ chính xác thấp hơn\nnhưng vẫn có kết quả", ACCENT_BLUE),
        ("🤖 LLM Failed", "→ Default tags", "Dùng tags mặc định\nthay vì crash", ACCENT_PURPLE),
        ("⏱️ Solver Timeout", "→ Best feasible solution", "Trả kết quả tốt nhất\ntrong time_limit", ACCENT_GREEN),
        ("🌐 Slow Request", "→ Rate limiting", "slowapi giới hạn request\ntránh quá tải", ACCENT_ORANGE),
    ]

    x = Inches(0.5)
    for trigger, action, detail, color in fallbacks:
        add_rounded_rect(slide, x, Inches(2.3), Inches(2.9), Inches(3.5), BG_CARD)
        add_accent_line(slide, x, Inches(2.3), Inches(2.9), color)
        add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(2.5), Inches(0.5),
                     trigger, font_size=15, bold=True, color=color)
        add_text_box(slide, x + Inches(0.2), Inches(3.1), Inches(2.5), Inches(0.4),
                     action, font_size=16, bold=True, color=ACCENT_GREEN)
        add_text_box(slide, x + Inches(0.2), Inches(3.7), Inches(2.5), Inches(0.8),
                     detail, font_size=12, color=GRAY_LIGHT)
        x += Inches(3.15)

    add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
                 "Goal: Ensure robust itinerary generation under all conditions",
                 font_size=16, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)


def build_slide_18_demo(prs):
    """Slide 18 — Demo Scenario"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Demo: Example Journey — 2 ngày Huế",
                      'Prompt: "2 ngày ở Huế, thích văn hóa, budget thấp"', 18)

    # Day 1
    add_rounded_rect(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(3.5), BG_CARD)
    add_accent_line(slide, Inches(0.8), Inches(2.3), Inches(5.5), ACCENT_BLUE)
    add_text_box(slide, Inches(1.1), Inches(2.5), Inches(5), Inches(0.4),
                 "📅 Ngày 1 — Văn hóa & Di tích", font_size=17, bold=True, color=ACCENT_BLUE)
    day1 = [
        "08:00  Đại Nội Huế (150k₫, 2h)",
        "10:30  Chùa Thiên Mụ (Free, 1h)",
        "12:00  Cơm Hến Bà Lý (30k₫, 45min)",
        "14:00  Lăng Tự Đức (100k₫, 1.5h)",
        "16:00  Cafe Muối Huế (25k₫, 45min)",
    ]
    dy = Inches(3.1)
    for item in day1:
        add_text_box(slide, Inches(1.3), dy, Inches(4.8), Inches(0.3),
                     f"▸ {item}", font_size=13, color=GRAY_LIGHT)
        dy += Inches(0.35)

    # Day 2
    add_rounded_rect(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(3.5), BG_CARD)
    add_accent_line(slide, Inches(7), Inches(2.3), Inches(5.5), ACCENT_GREEN)
    add_text_box(slide, Inches(7.3), Inches(2.5), Inches(5), Inches(0.4),
                 "📅 Ngày 2 — Thiên nhiên & Ẩm thực", font_size=17, bold=True, color=ACCENT_GREEN)
    day2 = [
        "08:00  Sông Hương boat (80k₫, 1.5h)",
        "10:00  Chợ Đông Ba (Free, 1h)",
        "11:30  Bún Bò Huế O Phụng (35k₫, 45min)",
        "13:30  Đồi Vọng Cảnh (Free, 1.5h)",
        "15:30  Cầu Trường Tiền (Free, 30min)",
    ]
    dy = Inches(3.1)
    for item in day2:
        add_text_box(slide, Inches(7.3), dy, Inches(4.8), Inches(0.3),
                     f"▸ {item}", font_size=13, color=GRAY_LIGHT)
        dy += Inches(0.35)

    # Summary
    add_rounded_rect(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.8),
                     RGBColor(0x0F, 0x1F, 0x15))
    add_text_box(slide, Inches(1.1), Inches(6.2), Inches(11), Inches(0.5),
                 "Total: 10 stops  •  18.5 km  •  420k₫  •  Optimized by OR-Tools CVRPTW",
                 font_size=16, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)


def build_slide_19_evaluation(prs):
    """Slide 19 — Evaluation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "System Evaluation",
                      "Metrics đánh giá hiệu suất hệ thống", 19)

    metrics = [
        ("Itinerary Generation", "~8-15s", "End-to-end pipeline\n(L2+L3+L4)", ACCENT_BLUE),
        ("Re-route Response", "~3-5s", "JIT single-day\nre-optimization", ACCENT_GREEN),
        ("Route Efficiency", "↓ 25-40%", "Travel distance\nvs manual planning", ACCENT_PURPLE),
        ("POI Relevance", "85%+", "Semantic match\nvs user intent", ACCENT_ORANGE),
    ]

    x = Inches(0.5)
    for title, value, desc, color in metrics:
        add_rounded_rect(slide, x, Inches(2.3), Inches(2.9), Inches(2.8), BG_CARD)
        add_accent_line(slide, x, Inches(2.3), Inches(2.9), color)
        add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(2.5), Inches(0.4),
                     title, font_size=14, bold=True, color=color)
        add_text_box(slide, x + Inches(0.2), Inches(3.0), Inches(2.5), Inches(0.7),
                     value, font_size=36, bold=True, color=WHITE)
        add_text_box(slide, x + Inches(0.2), Inches(3.8), Inches(2.5), Inches(0.8),
                     desc, font_size=12, color=GRAY_LIGHT)
        x += Inches(3.15)


def build_slide_20_contribution(prs):
    """Slide 20 — Contributions"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Key Contributions", "Đóng góp chính của đề tài", 20)

    contribs = [
        ("🤖", "NLP → Itinerary Extraction",
         "Tự động chuyển prompt tiếng Việt thành ràng buộc tối ưu hóa", ACCENT_BLUE),
        ("🔍", "Hybrid Semantic + Spatial Search",
         "Kết hợp vector similarity với geographic filtering trên PostGIS + pgvector", ACCENT_PURPLE),
        ("⚡", "Dynamic Multi-day Optimization",
         "2-stage CVRPTW solver: allocation + routing bằng OR-Tools", ACCENT_GREEN),
        ("🔄", "Real-time JIT Re-routing",
         "Virtual Depot re-optimization từ GPS hiện tại, response < 5s", ACCENT_ORANGE),
    ]

    y = Inches(2.3)
    for emoji, title, desc, color in contribs:
        add_rounded_rect(slide, Inches(0.8), y, Inches(11.5), Inches(1.05), BG_CARD)
        add_accent_line(slide, Inches(0.8), y, Inches(0.2), color)
        add_text_box(slide, Inches(1.3), y + Inches(0.1), Inches(10.5), Inches(0.4),
                     f"{emoji}  {title}", font_size=18, bold=True, color=color)
        add_text_box(slide, Inches(1.3), y + Inches(0.55), Inches(10.5), Inches(0.4),
                     desc, font_size=14, color=GRAY_LIGHT)
        y += Inches(1.2)


def build_slide_21_limitations(prs):
    """Slide 21 — Limitations & Future Work"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_block(slide, "Limitations & Future Work", None, 21)

    # LEFT — Limitations
    add_rounded_rect(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.8),
                     RGBColor(0x1E, 0x15, 0x15))
    add_text_box(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(0.4),
                 "⚠️ Current Limitations", font_size=18, bold=True, color=ACCENT_ORANGE)
    limitations = [
        "Phụ thuộc OSRM availability",
        "LLM có thể hallucinate intent",
        "Background location chưa auto-trigger",
        "Chưa hoàn thiện FCM production",
        "Cần thêm integration test xuyên layer",
    ]
    ly = Inches(2.8)
    for item in limitations:
        add_text_box(slide, Inches(1.3), ly, Inches(4.8), Inches(0.35),
                     f"• {item}", font_size=13, color=GRAY_LIGHT)
        ly += Inches(0.4)

    # RIGHT — Future
    add_rounded_rect(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(4.8),
                     RGBColor(0x0F, 0x1F, 0x15))
    add_text_box(slide, Inches(7.3), Inches(2.2), Inches(5), Inches(0.4),
                 "🚀 Future Improvements", font_size=18, bold=True, color=ACCENT_GREEN)
    future = [
        "Auto-trigger reroute via geofence/delay",
        "Weather-aware optimization",
        "Hotel/booking integration",
        "Crowd prediction & avoidance",
        "RL-based personalization",
        "Observability dashboard",
    ]
    fy = Inches(2.8)
    for item in future:
        add_text_box(slide, Inches(7.3), fy, Inches(4.8), Inches(0.35),
                     f"→ {item}", font_size=13, color=GRAY_LIGHT)
        fy += Inches(0.4)


def build_slide_22_qa(prs):
    """Slide 22 — Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Decorative lines
    add_accent_line(slide, Inches(4.5), Inches(2.5), Inches(4.5), ACCENT_BLUE)
    add_accent_line(slide, Inches(5.0), Inches(2.6), Inches(3.5), ACCENT_GREEN)

    # Big text
    add_text_box(slide, Inches(0), Inches(2.8), Inches(13.333), Inches(1.2),
                 "Thank You", font_size=56, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(4.0), Inches(13.333), Inches(0.8),
                 "Questions & Discussion", font_size=28, color=ACCENT_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # Suggested questions
    questions = [
        "Điểm khác biệt kỹ thuật lớn nhất so với app du lịch thông thường?",
        "Khi nào hệ thống dùng fallback và ảnh hưởng chất lượng route?",
        "Lộ trình chuyển từ POC sang production?",
    ]
    qy = Inches(5.2)
    for q in questions:
        add_text_box(slide, Inches(2), qy, Inches(9.5), Inches(0.35),
                     f"💬  {q}", font_size=13, color=GRAY_MID,
                     alignment=PP_ALIGN.CENTER)
        qy += Inches(0.4)
